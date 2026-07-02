from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import landscape
from reportlab.lib.styles import ParagraphStyle
from reportlab.lib.units import inch
from reportlab.pdfbase.pdfmetrics import stringWidth
from reportlab.pdfgen import canvas
from reportlab.platypus import Paragraph


ROOT = Path(__file__).resolve().parents[1]
ASSET_DIR = ROOT / "pitch-assets"
OUT_DIR = ROOT / "pitch"
OUT_DIR.mkdir(exist_ok=True)

TITLE_BG = ASSET_DIR / "title_bg.jpg"
BODY_BG = ASSET_DIR / "body_bg.jpg"
CLOSING_BG = ASSET_DIR / "closing_bg.jpg"

PDF_PATH = OUT_DIR / "CallOlve_AI_Hackathon_Pitch_Deck.pdf"
PPTX_PATH = OUT_DIR / "CallOlve_AI_Hackathon_Pitch_Deck.pptx"

W, H = 960, 540

NAVY = "#15133f"
INK = "#1f2d4d"
MUTED = "#667085"
PURPLE = "#7047d7"
BLUE = "#2258f5"
ORANGE = "#ff8a2a"
GREEN = "#20a66a"
LIGHT = "#f6f7fb"
WHITE = "#ffffff"
BLACK = "#000000"


def hex_rgb(value: str) -> tuple[int, int, int]:
    value = value.lstrip("#")
    return tuple(int(value[i : i + 2], 16) for i in (0, 2, 4))


def rgb(value: str) -> RGBColor:
    return RGBColor(*hex_rgb(value))


def prep_images() -> None:
    sources = {
        TITLE_BG: (ASSET_DIR / "image2.png", 78),
        BODY_BG: (ASSET_DIR / "image1.png", 90),
        CLOSING_BG: (ASSET_DIR / "image3.png", 78),
    }
    for dest, (src, quality) in sources.items():
        if dest.exists():
            continue
        im = Image.open(src).convert("RGB")
        im.save(dest, "JPEG", quality=quality, optimize=True, progressive=True)


slides = [
    {
        "kind": "title",
        "title": "CallOlve AI",
        "subtitle": "AI voice operators for every phone call",
        "body": [
            ("Team Name", "CallOlve AI"),
            ("Problem Statement", "Long wait times and scripted phone bots in companies and government services."),
            ("Team Leader", "Adarsh Kumar"),
        ],
    },
    {
        "title": "Problem: phone support still wastes caller time",
        "lead": "Whether it is a private company or a government service, callers often wait too long before reaching someone useful.",
        "cards": [
            ("Long queues", "People wait through ringing, hold music, and repeated transfers before reaching a representative."),
            ("Scripted IVR", "Recorded menus and basic bots follow fixed paths even when the caller has a real question."),
            ("No intelligence", "Most voice automation cannot reason, clarify, or adapt when the situation changes."),
            ("Lost trust", "Customers repeat details again and again, while teams lose context and operating time."),
        ],
    },
    {
        "title": "Solution overview",
        "lead": "CallOlve AI gives organizations an intelligent voice representative that can understand callers, respond naturally, and complete work.",
        "split": [
            ("What it does", [
                "Answers phone calls or browser voice sessions as a trained company agent.",
                "Understands caller intent, asks follow-up questions, and captures required details.",
                "Books appointments, takes orders, qualifies leads, handles support, and escalates to a human when needed.",
                "Stores transcript, summary, outcome, and next action for the company dashboard.",
            ]),
            ("Why it is different", [
                "Reasoning-first, not menu-first: the caller can speak normally instead of following a script.",
                "The agent uses context and company rules instead of a fixed recorded message.",
                "Provider-adapter design supports Twilio, LiveKit, Azure, ElevenLabs, and future vendors.",
                "Companies can rely on fallback, logs, and policy boundaries before scaling it live.",
            ]),
        ],
    },
    {
        "title": "What we built so far",
        "lead": "A working prototype that proves the core call-agent flow, while the final human-grade voice experience is still the funded roadmap.",
        "modules": [
            ("Assistant builder", "Roles, personality, tone, voice, language, greeting, prompt, memory."),
            ("Live voice", "Browser voice path plus phone worker using Twilio SIP, LiveKit, and Azure realtime."),
            ("Action engine", "Tool calls for bookings, orders, leads, human handoff, and callbacks."),
            ("Company dashboard", "Calls, assistants, appointments, orders, leads, analytics, integrations."),
            ("Call memory", "Transcript, summary, caller intent, extracted details, and next action."),
            ("Persistence", "Shared API payload for simulator and phone calls, backed by Prisma data models."),
        ],
    },
    {
        "title": "Why we built it this way",
        "lead": "The architecture is designed to move from a hackathon prototype into a reliable company phone representative.",
        "rationale": [
            ("One contract", "Simulator, browser voice, and PSTN calls all produce the same structured Call payload."),
            ("Real-time media layer", "LiveKit handles room media and SIP dispatch; the worker connects outbound, so no public tunnel is required."),
            ("Realtime model path", "Azure Foundry gpt-realtime-2 removes the separate STT -> LLM -> TTS loop for phone calls."),
            ("Natural turn-taking", "The roadmap targets interruptible, ultra-responsive voice that can listen and reason like a human representative."),
            ("Business reliability", "The assistant confirms details before taking action and hands off to humans when confidence is low."),
        ],
    },
    {
        "title": "End-to-end workflow",
        "lead": "From a ringing phone to an intelligent answer, action, or human handoff.",
        "workflow": [
            ("1", "Caller calls", "The company number receives the call."),
            ("2", "Media room", "LiveKit connects the voice stream."),
            ("3", "AI agent", "Worker loads company persona and rules."),
            ("4", "Reason + act", "Agent clarifies, answers, or executes."),
            ("5", "Record", "Transcript and next step are saved."),
        ],
    },
    {
        "title": "System architecture",
        "lead": "A modular voice platform with clear ownership boundaries.",
        "architecture": True,
    },
    {
        "title": "Intelligent agent behavior",
        "lead": "CallOlve is designed to replace rigid scripts with a company-trained agent that can reason during the call.",
        "cards": [
            ("Persona rendering", "Each assistant prompt is generated from role, tone, language, greeting, and custom instructions."),
            ("Structured tools", "Booking, order, lead, handoff, and callback tools return deterministic action payloads."),
            ("Anti-hallucination", "The assistant only claims completion after a tool executes and logs the resulting action."),
            ("Data quality", "Names, times, items, phone context, and caller intent are collected as slots before persistence."),
            ("Human handoff", "When confidence is low or policy requires it, the agent routes the caller to a real representative."),
            ("Audit trail", "Every completed call keeps transcript, outcome, action list, and dashboard record."),
        ],
    },
    {
        "title": "Current status and roadmap",
        "lead": "CallOlve AI is not yet the complete production product. The prototype validates the core voice-agent path; funding unlocks the human-grade version.",
        "metrics": [
            ("LiveKit SIP dispatch", "Verified inbound trunk and dispatch rule route calls to callolve-phone."),
            ("Azure realtime", "gpt-realtime-2 websocket accepted the session config and runs as the phone LLM/voice path."),
            ("Call persistence", "Test SIP calls generated assistant greetings and were saved through /api/v1/voice/complete."),
            ("Funding roadmap", "Human-realistic voice, lower latency, stronger barge-in, and agents that listen while preparing the next response."),
        ],
        "demo": [
            "Demo flow: call the number -> Aria greets -> caller requests booking/order -> assistant confirms -> dashboard record appears.",
            "Target product: a real-time voice representative that companies and public services can trust at scale.",
        ],
    },
    {
        "title": "Technologies used",
        "lead": "Chosen for a fast hackathon build that still maps to production deployment.",
        "tech": [
            ("Frontend", "Next.js App Router, React, TypeScript, Tailwind-style design tokens."),
            ("Backend", "Next.js API routes, Prisma, SQLite for local demo, structured service layer."),
            ("Voice", "Twilio PSTN/SIP, LiveKit Cloud + LiveKit Agents, Azure Foundry gpt-realtime-2."),
            ("AI/tools", "OpenAI-compatible LLM path, function tools, shared action schema, prompt renderer."),
            ("Voice roadmap", "Human-realistic TTS, interruptible speech, faster streaming, and parallel listening/reasoning."),
            ("Validation", "Preflight checks for provider config, SIP setup, realtime websocket, and worker logs."),
        ],
    },
    {
        "title": "Submission assets and next steps",
        "lead": "What to submit now, and what funding would turn into the full company-ready product.",
        "assets": [
            ("PDF deck", "This file: concise approach, build scope, architecture, and validation."),
            ("Demo video", "Recommended: 90 seconds showing assistant setup, phone call, and persisted dashboard record."),
            ("Repository", "Add the final GitHub URL in the H2S dashboard when the project is ready to publish."),
            ("Next steps", "Fund human-like voice, ultra-low latency, stronger reliability, and production integrations."),
        ],
    },
    {
        "kind": "closing",
        "title": "CallOlve AI",
        "subtitle": "The operating system for voice communication.",
        "footer": "Built for India.Runs - AI calls that actually get work done.",
    },
]


def wrap_lines(text: str, font: str, size: float, max_width: float) -> list[str]:
    words = text.split()
    lines: list[str] = []
    cur = ""
    for word in words:
        trial = f"{cur} {word}".strip()
        if stringWidth(trial, font, size) <= max_width:
            cur = trial
        else:
            if cur:
                lines.append(cur)
            cur = word
    if cur:
        lines.append(cur)
    return lines


def pdf_text(c: canvas.Canvas, text: str, x: float, y: float, size: float, color: str = INK, bold: bool = False, max_width: float | None = None, leading: float | None = None) -> float:
    font = "Helvetica-Bold" if bold else "Helvetica"
    c.setFont(font, size)
    c.setFillColor(colors.HexColor(color))
    leading = leading or size * 1.25
    if max_width is None:
        c.drawString(x, y, text)
        return y - leading
    for line in wrap_lines(text, font, size, max_width):
        c.drawString(x, y, line)
        y -= leading
    return y


def pdf_bg(c: canvas.Canvas, kind: str = "body") -> None:
    img = BODY_BG
    if kind == "title":
        img = TITLE_BG
    elif kind == "closing":
        img = CLOSING_BG
    c.drawImage(str(img), 0, 0, width=W, height=H)


def pdf_card(c: canvas.Canvas, x: float, y: float, w: float, h: float, title: str, body: str, accent: str = PURPLE) -> None:
    c.setFillColor(colors.HexColor(WHITE))
    c.setStrokeColor(colors.HexColor("#d9dde8"))
    c.roundRect(x, y, w, h, 10, fill=1, stroke=1)
    c.setFillColor(colors.HexColor(accent))
    c.roundRect(x, y + h - 9, w, 9, 6, fill=1, stroke=0)
    pdf_text(c, title, x + 16, y + h - 34, 15, INK, True, w - 32)
    pdf_text(c, body, x + 16, y + h - 62, 10.5, MUTED, False, w - 32, 14)


def pdf_bullets(c: canvas.Canvas, items: list[str], x: float, y: float, w: float, size: float = 11.5) -> float:
    for item in items:
        c.setFillColor(colors.HexColor(PURPLE))
        c.circle(x + 4, y + 4, 3, fill=1, stroke=0)
        y = pdf_text(c, item, x + 16, y, size, INK, False, w - 16, size * 1.25) - 5
    return y


def pdf_header(c: canvas.Canvas, slide: dict) -> float:
    pdf_bg(c, "body")
    pdf_text(c, slide["title"], 44, 431, 27, NAVY, True, 840, 34)
    return pdf_text(c, slide.get("lead", ""), 44, 392, 13, MUTED, False, 820, 17)


def draw_pdf() -> None:
    c = canvas.Canvas(str(PDF_PATH), pagesize=landscape((W, H)))

    for idx, slide in enumerate(slides, start=1):
        kind = slide.get("kind")
        if kind == "title":
            pdf_bg(c, "title")
            c.setFillColor(colors.HexColor("#10112a"))
            c.roundRect(42, 118, 850, 235, 16, fill=1, stroke=0)
            pdf_text(c, slide["title"], 55, 306, 54, WHITE, True)
            pdf_text(c, slide["subtitle"], 58, 266, 22, "#e8eaff", False)
            y = 190
            for label, value in slide["body"]:
                pdf_text(c, f"{label}:", 58, y, 13, "#d7d8ff", True)
                y = pdf_text(c, value, 190, y, 13, WHITE, False, 690, 17) - 8
            c.showPage()
            continue

        if kind == "closing":
            pdf_bg(c, "closing")
            c.showPage()
            continue

        y = pdf_header(c, slide)

        if "cards" in slide:
            cards = slide["cards"]
            cols = 2 if len(cards) <= 4 else 3
            cw = 410 if cols == 2 else 270
            ch = 105 if len(cards) <= 4 else 88
            start_x = 44
            start_y = 245 if len(cards) <= 4 else 268
            for i, (title, body) in enumerate(cards):
                x = start_x + (i % cols) * (cw + 24)
                yy = start_y - (i // cols) * (ch + 22)
                pdf_card(c, x, yy, cw, ch, title, body, [PURPLE, BLUE, ORANGE, GREEN, PURPLE, BLUE][i % 6])

        elif "split" in slide:
            for col, (title, bullets) in enumerate(slide["split"]):
                x = 52 + col * 438
                c.setFillColor(colors.HexColor(LIGHT))
                c.roundRect(x, 78, 385, 275, 12, fill=1, stroke=0)
                pdf_text(c, title, x + 22, 320, 18, INK, True)
                pdf_bullets(c, bullets, x + 24, 286, 335, 11)

        elif "modules" in slide:
            for i, (title, body) in enumerate(slide["modules"]):
                x = 52 + (i % 3) * 292
                yy = 276 - (i // 3) * 120
                pdf_card(c, x, yy, 255, 94, title, body, [PURPLE, BLUE, ORANGE, GREEN, PURPLE, BLUE][i])

        elif "rationale" in slide:
            y2 = 330
            for i, (title, body) in enumerate(slide["rationale"]):
                c.setFillColor(colors.HexColor([PURPLE, BLUE, ORANGE, GREEN, PURPLE][i]))
                c.circle(64, y2 + 3, 10, fill=1, stroke=0)
                c.setFillColor(colors.white)
                c.setFont("Helvetica-Bold", 9)
                c.drawCentredString(64, y2, str(i + 1))
                pdf_text(c, title, 88, y2 + 8, 15, INK, True)
                pdf_text(c, body, 88, y2 - 10, 11, MUTED, False, 785, 14)
                y2 -= 56

        elif "workflow" in slide:
            x = 44
            y2 = 253
            for i, (num, title, body) in enumerate(slide["workflow"]):
                c.setFillColor(colors.HexColor(WHITE))
                c.setStrokeColor(colors.HexColor("#d9dde8"))
                c.roundRect(x, y2, 160, 108, 12, fill=1, stroke=1)
                c.setFillColor(colors.HexColor([PURPLE, BLUE, ORANGE, GREEN, PURPLE][i]))
                c.circle(x + 24, y2 + 78, 14, fill=1, stroke=0)
                c.setFillColor(colors.white)
                c.setFont("Helvetica-Bold", 12)
                c.drawCentredString(x + 24, y2 + 74, num)
                pdf_text(c, title, x + 18, y2 + 52, 13, INK, True, 128)
                pdf_text(c, body, x + 18, y2 + 30, 9.5, MUTED, False, 128, 12)
                if i < len(slide["workflow"]) - 1:
                    c.setStrokeColor(colors.HexColor(PURPLE))
                    c.setLineWidth(1.6)
                    c.line(x + 165, y2 + 54, x + 190, y2 + 54)
                    c.line(x + 184, y2 + 60, x + 190, y2 + 54)
                    c.line(x + 184, y2 + 48, x + 190, y2 + 54)
                x += 184

        elif slide.get("architecture"):
            boxes = [
                (52, 292, 185, 70, "Caller channels", "Phone, browser voice, web app"),
                (282, 292, 185, 70, "Voice transport", "Twilio SIP + LiveKit rooms"),
                (512, 292, 185, 70, "AI worker", "LiveKit Agents + Azure realtime"),
                (742, 292, 150, 70, "Actions", "Book, order, lead, handoff"),
                (170, 138, 210, 72, "CallOlve API", "Next.js routes + shared contract"),
                (430, 138, 210, 72, "Data layer", "Prisma, transcripts, slots"),
                (690, 138, 190, 72, "Dashboard", "Analytics, history, follow-up"),
            ]
            for i, (x, yy, ww, hh, title, body) in enumerate(boxes):
                pdf_card(c, x, yy, ww, hh, title, body, [PURPLE, BLUE, ORANGE, GREEN, PURPLE, BLUE, GREEN][i])
            for x1, y1, x2, y2 in [(237, 327, 282, 327), (467, 327, 512, 327), (697, 327, 742, 327), (585, 292, 535, 210), (742, 292, 785, 210), (380, 174, 430, 174), (640, 174, 690, 174)]:
                c.setStrokeColor(colors.HexColor("#6e6aa8"))
                c.setLineWidth(1.6)
                c.line(x1, y1, x2, y2)

        elif "metrics" in slide:
            for i, (title, body) in enumerate(slide["metrics"]):
                x = 55 + (i % 2) * 420
                yy = 260 - (i // 2) * 102
                pdf_card(c, x, yy, 375, 78, title, body, [GREEN, BLUE, PURPLE, ORANGE][i])
            pdf_text(c, "Demo proof points", 56, 102, 16, INK, True)
            pdf_bullets(c, slide["demo"], 58, 78, 810, 11)

        elif "tech" in slide:
            y2 = 336
            for i, (title, body) in enumerate(slide["tech"]):
                x = 54 + (i % 2) * 430
                yy = y2 - (i // 2) * 70
                c.setFillColor(colors.HexColor("#f8f9fe"))
                c.roundRect(x, yy - 40, 375, 55, 8, fill=1, stroke=0)
                pdf_text(c, title, x + 16, yy, 12.5, [PURPLE, BLUE, ORANGE, GREEN, PURPLE, BLUE][i], True)
                pdf_text(c, body, x + 116, yy, 10, INK, False, 235, 12)

        elif "assets" in slide:
            for i, (title, body) in enumerate(slide["assets"]):
                x = 80
                yy = 315 - i * 64
                c.setFillColor(colors.HexColor([PURPLE, BLUE, ORANGE, GREEN][i]))
                c.roundRect(x, yy - 11, 95, 34, 15, fill=1, stroke=0)
                c.setFillColor(colors.white)
                c.setFont("Helvetica-Bold", 11)
                c.drawCentredString(x + 47, yy, title)
                pdf_text(c, body, x + 122, yy + 4, 12, INK, False, 665, 15)

        pdf_text(c, f"{idx:02d}", 886, 31, 10, "#80839a", True)
        c.showPage()

    c.save()


def ppt_textbox(slide, x, y, w, h, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color)
    return tb


def ppt_bullets(slide, bullets: list[str], x, y, w, h, size=12, color=INK):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(color)
    return tb


def ppt_card(slide, x, y, w, h, title, body, accent=PURPLE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(WHITE)
    shape.line.color.rgb = rgb("#d9dde8")
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(accent)
    bar.line.fill.background()
    ppt_textbox(slide, x + 0.16, y + 0.19, w - 0.32, 0.28, title, 10.5, INK, True)
    ppt_textbox(slide, x + 0.16, y + 0.50, w - 0.32, h - 0.52, body, 8.2, MUTED)


def ppt_background(slide, kind="body"):
    img = BODY_BG
    if kind == "title":
        img = TITLE_BG
    elif kind == "closing":
        img = CLOSING_BG
    slide.shapes.add_picture(str(img), 0, 0, width=Inches(10), height=Inches(5.625))


def ppt_header(slide, item):
    ppt_background(slide)
    ppt_textbox(slide, 0.48, 0.89, 9.0, 0.45, item["title"], 22, NAVY, True)
    ppt_textbox(slide, 0.48, 1.33, 8.75, 0.42, item.get("lead", ""), 10.5, MUTED)


def draw_pptx() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    blank = prs.slide_layouts[6]

    for idx, item in enumerate(slides, start=1):
        slide = prs.slides.add_slide(blank)
        kind = item.get("kind")
        if kind == "title":
            ppt_background(slide, "title")
            panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.42), Inches(1.88), Inches(9.05), Inches(2.95))
            panel.fill.solid()
            panel.fill.fore_color.rgb = rgb("#10112a")
            panel.line.fill.background()
            ppt_textbox(slide, 0.58, 2.08, 7.0, 0.62, item["title"], 42, WHITE, True)
            ppt_textbox(slide, 0.62, 2.72, 7.5, 0.35, item["subtitle"], 17, "#e8eaff")
            y = 3.55
            for label, value in item["body"]:
                ppt_textbox(slide, 0.62, y, 1.35, 0.22, label + ":", 9.5, "#d7d8ff", True)
                ppt_textbox(slide, 2.03, y, 7.25, 0.33, value, 9.5, WHITE)
                y += 0.43
            continue
        if kind == "closing":
            ppt_background(slide, "closing")
            continue

        ppt_header(slide, item)

        if "cards" in item:
            cards = item["cards"]
            cols = 2 if len(cards) <= 4 else 3
            cw = 4.25 if cols == 2 else 2.78
            ch = 1.05 if len(cards) <= 4 else 0.92
            start_x = 0.55
            start_y = 2.63 if len(cards) <= 4 else 2.72
            for i, (title, body) in enumerate(cards):
                x = start_x + (i % cols) * (cw + 0.30)
                y = start_y + (i // cols) * (ch + 0.26)
                ppt_card(slide, x, y, cw, ch, title, body, [PURPLE, BLUE, ORANGE, GREEN, PURPLE, BLUE][i % 6])
        elif "split" in item:
            for col, (title, bullets) in enumerate(item["split"]):
                x = 0.62 + col * 4.55
                box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.95), Inches(4.05), Inches(2.85))
                box.fill.solid()
                box.fill.fore_color.rgb = rgb(LIGHT)
                box.line.fill.background()
                ppt_textbox(slide, x + 0.22, 2.16, 3.5, 0.28, title, 15, INK, True)
                ppt_bullets(slide, bullets, x + 0.28, 2.55, 3.4, 1.85, 9.5)
        elif "modules" in item:
            for i, (title, body) in enumerate(item["modules"]):
                x = 0.58 + (i % 3) * 3.03
                y = 2.08 + (i // 3) * 1.28
                ppt_card(slide, x, y, 2.66, 0.98, title, body, [PURPLE, BLUE, ORANGE, GREEN, PURPLE, BLUE][i])
        elif "rationale" in item:
            y = 2.0
            for i, (title, body) in enumerate(item["rationale"]):
                circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.62), Inches(y), Inches(0.26), Inches(0.26))
                circ.fill.solid()
                circ.fill.fore_color.rgb = rgb([PURPLE, BLUE, ORANGE, GREEN, PURPLE][i])
                circ.line.fill.background()
                ppt_textbox(slide, 0.96, y - 0.02, 2.2, 0.25, title, 12.5, INK, True)
                ppt_textbox(slide, 0.96, y + 0.25, 8.45, 0.32, body, 9.5, MUTED)
                y += 0.62
        elif "workflow" in item:
            for i, (num, title, body) in enumerate(item["workflow"]):
                x = 0.48 + i * 1.92
                y = 2.42
                ppt_card(slide, x, y, 1.66, 1.25, title, body, [PURPLE, BLUE, ORANGE, GREEN, PURPLE][i])
                circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.06), Inches(y - 0.16), Inches(0.28), Inches(0.28))
                circ.fill.solid()
                circ.fill.fore_color.rgb = rgb([PURPLE, BLUE, ORANGE, GREEN, PURPLE][i])
                circ.line.fill.background()
                ppt_textbox(slide, x + 0.13, y - 0.11, 0.14, 0.12, num, 7.5, WHITE, True, PP_ALIGN.CENTER)
                if i < len(item["workflow"]) - 1:
                    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + 1.69), Inches(y + 0.63), Inches(x + 1.88), Inches(y + 0.63))
                    line.line.color.rgb = rgb(PURPLE)
        elif item.get("architecture"):
            boxes = [
                (0.55, 2.05, 1.9, 0.86, "Caller channels", "Phone, browser voice, web app"),
                (2.93, 2.05, 1.9, 0.86, "Voice transport", "Twilio SIP + LiveKit rooms"),
                (5.31, 2.05, 1.9, 0.86, "AI worker", "LiveKit Agents + Azure realtime"),
                (7.69, 2.05, 1.55, 0.86, "Actions", "Book, order, lead"),
                (1.75, 3.55, 2.16, 0.86, "CallOlve API", "Next.js routes + shared contract"),
                (4.45, 3.55, 2.16, 0.86, "Data layer", "Prisma, transcripts, slots"),
                (7.15, 3.55, 1.96, 0.86, "Dashboard", "Analytics, history, follow-up"),
            ]
            for i, (x, y, w, h, title, body) in enumerate(boxes):
                ppt_card(slide, x, y, w, h, title, body, [PURPLE, BLUE, ORANGE, GREEN, PURPLE, BLUE, GREEN][i])
        elif "metrics" in item:
            for i, (title, body) in enumerate(item["metrics"]):
                x = 0.58 + (i % 2) * 4.35
                y = 2.08 + (i // 2) * 1.05
                ppt_card(slide, x, y, 3.9, 0.85, title, body, [GREEN, BLUE, PURPLE, ORANGE][i])
            ppt_textbox(slide, 0.62, 4.18, 2.3, 0.25, "Demo proof points", 13, INK, True)
            ppt_bullets(slide, item["demo"], 0.62, 4.48, 8.4, 0.55, 9.5)
        elif "tech" in item:
            for i, (title, body) in enumerate(item["tech"]):
                x = 0.58 + (i % 2) * 4.45
                y = 1.95 + (i // 2) * 0.74
                box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.92), Inches(0.55))
                box.fill.solid()
                box.fill.fore_color.rgb = rgb("#f8f9fe")
                box.line.fill.background()
                ppt_textbox(slide, x + 0.15, y + 0.10, 1.05, 0.2, title, 9.5, [PURPLE, BLUE, ORANGE, GREEN, PURPLE, BLUE][i], True)
                ppt_textbox(slide, x + 1.18, y + 0.08, 2.5, 0.32, body, 8.3, INK)
        elif "assets" in item:
            for i, (title, body) in enumerate(item["assets"]):
                y = 2.02 + i * 0.68
                pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.86), Inches(y), Inches(1.1), Inches(0.34))
                pill.fill.solid()
                pill.fill.fore_color.rgb = rgb([PURPLE, BLUE, ORANGE, GREEN][i])
                pill.line.fill.background()
                ppt_textbox(slide, 0.93, y + 0.06, 0.95, 0.14, title, 8.5, WHITE, True, PP_ALIGN.CENTER)
                ppt_textbox(slide, 2.16, y + 0.03, 7.1, 0.3, body, 10.2, INK)

        ppt_textbox(slide, 9.25, 5.1, 0.45, 0.18, f"{idx:02d}", 7, "#80839a", True, PP_ALIGN.RIGHT)

    prs.save(PPTX_PATH)


if __name__ == "__main__":
    prep_images()
    draw_pptx()
    draw_pdf()
    print(PPTX_PATH)
    print(PDF_PATH)
